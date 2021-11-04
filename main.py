from pptx import Presentation
from pptx.util import Pt
import os
import pandas as pd

# import comtypes.client
import smtplib, ssl

from email import encoders
from email.mime.base import MIMEBase
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText

from dotenv import load_dotenv

load_dotenv()


def delete_file(filename):
    try:
        os.remove(filename)
    except Exception as e:
        print(e)


def delete_pptx(names):
    for name in names:
        try:
            os.remove(f"output/Output_{name}.pptx")
        except Exception as e:
            print(e)
            pass


def read_csv(filename):
    df = pd.read_csv(filename)
    emails = df.iloc[:, 1]
    names = df.iloc[:, 0]
    return names, emails


def add_name_to_ppt(ppt_file_path, names):
    # Opening file
    for name in names:
        prs = Presentation(ppt_file_path)
        slide = prs.slides[0]

        # Prinitng everyones name on the certificate
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame

                if text_frame.text == "Name":
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = name
                    font = run.font
                    font.name = "Caveat"
                    font.size = Pt(35)

        file_path_ppt = f"output/Output_{name}.pptx"
        file_path_pdf = f"output/Output_{name}.pdf"
        prs.save(file_path_ppt)
        # ppt_to_pdf(file_path_ppt, file_path_pdf)
        # delete_file(file_path_ppt)


def ppt_to_pdf(input_file_path, output_file_path):
    # Create powerpoint application object
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    # Set visibility to minimize
    powerpoint.Visible = 1
    # Open the powerpoint slides
    slides = powerpoint.Presentations.Open(input_file_path)
    # Save as PDF (formatType = 32)
    slides.SaveAs(output_file_path, 32)
    # Close the slide deck
    slides.Close()


def create_email_body(name, email):
    subject = "Google Cloud Completion Mail"
    body = "This is an email with attachment sent from Python"
    sender_email = os.environ.get("EMAIL")
    receiver_email = email

    # Create a multipart message and set headers
    message = MIMEMultipart()
    message["From"] = sender_email
    message["To"] = receiver_email
    message["Subject"] = subject
    message["Bcc"] = receiver_email  # Recommended for mass emails

    # Add body to email
    message.attach(MIMEText(body, "plain"))

    filename = f"output/Output_{name}.pptx" # In same directory as script
    # Open PDF file in binary mode
    with open(filename, "rb") as attachment:
        # Add file as application/octet-stream
        # Email client can usually download this automatically as attachment
        part = MIMEBase("application", "octet-stream")
        part.set_payload(attachment.read())

    # Encode file in ASCII characters to send by email
    encoders.encode_base64(part)

    # Add header as key/value pair to attachment part
    part.add_header(
        "Content-Disposition",
        f"attachment; filename= {name} Certificate.pptx",
    )

    # Add attachment to message and convert message to string
    message.attach(part)
    text = message.as_string()
    return text


def send_mail(names, emails):
    sender_email = os.environ.get("EMAIL")
    receiver_email = "akhilbhalerao@gmail.com"
    password = os.environ.get("PASSWORD")

    # Log in to server using secure context and send email
    context = ssl.create_default_context()
    with smtplib.SMTP_SSL("smtp.gmail.com", 465, context=context) as server:
        print("Logging In...")
        server.login(sender_email, password)
        print("Logged In...")
        print("Sending mail...")

        for i in range(len(emails)):
            text = create_email_body(names[i], emails[i])
            server.sendmail(sender_email, receiver_email, text)
            print(f"Mail sent successfully to {names[i]}.")


if __name__ == "__main__":
    user_name = "Akhil Bhalerao"
    ppt_file_path = "assets/Both_tracks.pptx"
    user_data_file_path = os.path.abspath("User_details/Trial.csv")
    names, emails = read_csv(user_data_file_path)
    add_name_to_ppt(ppt_file_path, names)
    send_mail(names, emails)

    delete_pptx(names)
